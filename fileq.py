import streamlit as st
import openai
from pptx import Presentation
import os

openai.api_key = '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'


# Function to extract text from PPT
def extract_text_from_ppt(file):
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs


# Function to generate flashcards using OpenAI
def generate_flashcards(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo-2024-04-09",  # Choose the appropriate model
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"Generate a flashcard question and answer based on this content: {text}"}
            ],
            max_tokens=150,
            n=1,
            temperature=0.5
        )
        return response['choices'][0]['message']['content'].strip()
    except Exception as e:
        return str(e)

# Streamlit app
st.title('Flashcard Maker from PowerPoint Presentations')

uploaded_file = st.file_uploader("Upload your presentation", type=["pptx"])
if uploaded_file is not None:
    # Extract text from uploaded PPT
    extracted_texts = extract_text_from_ppt(uploaded_file)
    if st.button('Generate Flashcards'):
        # Display flashcards
        st.subheader("Generated Flashcards:")
        for text in extracted_texts:
            if text.strip():  # Check if the text is not just whitespace
                flashcard = generate_flashcards(text)
                st.text_area("Flashcard", value=flashcard, height=150)
